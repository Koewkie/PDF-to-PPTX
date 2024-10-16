import os, sys


from pdf2image import convert_from_path
from pptx import Presentation
from io import BytesIO

pdf_file = input("Enter the file name (must be in same directory)\n")
if not os.path.exists(pdf_file):	#check if file exists
	print("File not found")
	exit()

print()
print("Converting file: " + pdf_file)
print()

# Prep presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
# Create working folder
base_name = pdf_file.split(".pdf")[0]
# Convert PDF to list of images
print("Starting conversion...")
#Set poppler location
with open("Config.txt", "r") as cfgFile:
	popplerDir = cfgFile.readline()
slideimgs = convert_from_path(pdf_file,poppler_path=popplerDir)

# Loop over slides
for i, slideimg in enumerate(slideimgs):

	print("Saving slide: " + str(i))

	imagefile = BytesIO()	#create BytesIO object, dont have to write image to disk
	slideimg.save(imagefile, format='tiff')	#save slideimg to imagefile
	width, height = slideimg.size

	# Set slide dimensions
	prs.slide_height = height * 9525
	prs.slide_width = width * 9525

	# Add slide
	slide = prs.slides.add_slide(blank_slide_layout)
	pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

# Save Powerpoint
print()
print("Saving file: " + base_name + ".pptx")
prs.save(base_name + '.pptx')
print("Conversion complete. :)")
print()